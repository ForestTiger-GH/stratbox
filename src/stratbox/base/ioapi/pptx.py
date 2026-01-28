"""
pptx — чтение/запись PPTX поверх FileStore.

Важно:
- .ppt (старый бинарный формат) намеренно не разбирается: его лучше хранить как bytes.
- Для .pptx используется python-pptx (обычно установлен не всегда).

Сценарии:
- read_text(): извлечь весь текст со всех слайдов (для поиска/быстрого анализа)
- write_text(): создать простую презентацию: каждый элемент списка = отдельный слайд с заголовком
"""

from __future__ import annotations

from io import BytesIO
from typing import Iterable

from stratbox.base.filestore.base import FileStore
from stratbox.base.ioapi.bytes import read_bytes, write_bytes


def read_text(path: str, store: FileStore | None = None) -> str:
    """Читает PPTX и возвращает весь текст (по слайдам)."""
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise ImportError(
            "PPTX support requires optional dependency 'python-pptx'. "
            "Install: pip install python-pptx"
        ) from e

    data = read_bytes(path, store=store)
    bio = BytesIO(data)
    pres = Presentation(bio)

    out_lines: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_lines.append(shape.text)
        if slide_lines:
            out_lines.append(f"--- SLIDE {i} ---")
            out_lines.extend(slide_lines)

    return "\n".join(out_lines)


def write_text(path: str, titles: Iterable[str], store: FileStore | None = None) -> None:
    """
    Создаёт простую презентацию: каждый заголовок -> отдельный слайд.

    Это заготовка для быстрых прототипов.
    """
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise ImportError(
            "PPTX support requires optional dependency 'python-pptx'. "
            "Install: pip install python-pptx"
        ) from e

    pres = Presentation()
    layout = pres.slide_layouts[0]  # Title Slide

    for t in titles:
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = str(t)

    out = BytesIO()
    pres.save(out)
    write_bytes(path, out.getvalue(), store=store)
